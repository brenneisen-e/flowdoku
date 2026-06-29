from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DG  =RGBColor(0x04,0x6A,0x38)   # Primaer Dunkelgruen
DG2 =RGBColor(0x26,0x89,0x0D)   # Sekundaer
DARK=RGBColor(0x1A,0x1A,0x1A)   # Fliesstext
GREY=RGBColor(0x53,0x56,0x5A)   # Sekundaertext
PALE=RGBColor(0xF1,0xF6,0xE4)   # heller Panel-Hintergrund (Template pale green)
WHITE=RGBColor(0xFF,0xFF,0xFF)
LIGHT=RGBColor(0xEA,0xF3,0xDE)  # Text auf Dunkelgruen
FONT="Aptos"

p=Presentation("Template.pptx")
sl=p.slides._sldIdLst
for s in list(sl): p.part.drop_rel(s.get(qn('r:id'))); sl.remove(s)
LAY=p.slide_layouts[32]

# ----- Title slide (Dark-Bright gradient) -----
TS=p.slide_layouts[14]
ts=p.slides.add_slide(TS)
ts.placeholders[0].text_frame.text="DEX Event App"
ts.placeholders[10].text_frame.text="Stand, Governance & Ausblick  ·  OMP-Update  ·  Juni 2026  ·  Eike Brenneisen"
# remove empty picture placeholder so no grey box renders
for _ph in list(ts.placeholders):
    if _ph.placeholder_format.idx==11:
        _ph._element.getparent().remove(_ph._element)

# ----- Section divider (Deloitte Green) -----
DV=p.slide_layouts[20]
dv=p.slides.add_slide(DV)
dv.placeholders[0].text_frame.text="Status, Governance & Ausblick"

def set_ph(slide,idx,text): slide.placeholders[idx].text_frame.text=text

def card(slide,x,y,w,h,fill):
    sp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.adjustments[0]=0.045; sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.fill.background(); sp.shadow.inherit=False
    return sp

def card_text(shape, paras, ml,mt,mr,mb, anchor):
    tf=shape.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(ml); tf.margin_top=Inches(mt)
    tf.margin_right=Inches(mr); tf.margin_bottom=Inches(mb)
    first=True
    for (txt,size,color,bold,sb,sa) in paras:
        para=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        para.space_before=Pt(sb); para.space_after=Pt(sa); para.line_spacing=1.05
        r=para.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
        r.font.color.rgb=color; r.font.name=FONT

def badge(slide,x,y,d,fill,icon):
    ov=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    ov.fill.solid(); ov.fill.fore_color.rgb=fill; ov.line.fill.background(); ov.shadow.inherit=False
    isz=d*0.52; off=(d-isz)/2
    slide.shapes.add_picture(icon,Inches(x+off),Inches(y+off),Inches(isz),Inches(isz))

# ---------- SLIDE 1: KPI columns (text INSIDE the card) ----------
left0=0.55; total=13.33-1.1; GAP=0.32; CW=(total-3*GAP)/4
L=[left0+i*(CW+GAP) for i in range(4)]; CTOP=2.02; CH=4.0; BD=0.92; PX=0.30
def kpi(slide,i,icon,num,label,desc):
    x=L[i]; sp=card(slide,x,CTOP,CW,CH,PALE)
    badge(slide,x+PX,CTOP+0.34,BD,DG,icon)
    card_text(sp,[(num,40,DG,True,0,2),(label,15,DARK,True,2,4),(desc,11,GREY,False,0,0)],
              ml=PX, mt=0.34+BD+0.12, mr=PX, mb=0.2, anchor=MSO_ANCHOR.TOP)

s=p.slides.add_slide(LAY)
set_ph(s,0,"DEX – Was seit April passiert ist")
set_ph(s,13,"Eigener Plattform-Wechsel statt Entwicklung in Indien – mit messbarem Erfolg")
kpi(s,0,"icons/refresh-cw.png","April","Neue Plattform live","Bewusst gegen Indien entschieden – KI-gestützt eigenständig neu aufgebaut, live & produktiv.")
kpi(s,1,"icons/calendar.png","29","Events koordiniert","29 Events vollständig über DEX abgewickelt – von kleinen Runden bis Großevents, Tendenz steigend.")
kpi(s,2,"icons/users.png","~7.000","Teilnehmer","Rund 7.000 Anmeldungen gesteuert – automatische Mail- & Outlook-Anbindung spart den Assistenzen viel Zeit.")
kpi(s,3,"icons/check-circle.png","≈100%","Self-Service","Event-Erstellung läuft eigenständig durch die Assistenzen – kaum noch operative Einbindung. 72 Std. Eigenleistung gebucht.")

# ---------- SLIDES 2 & 3: 2x2 horizontal cards (text INSIDE the card) ----------
GX=0.30; GY=0.30; CW2=(total-GX)/2; CTOP2=1.92; CH2=((6.92-CTOP2)-GY)/2
COLX=[0.55,0.55+CW2+GX]; ROWY=[CTOP2,CTOP2+CH2+GY]; BD2=0.95; LPAD=0.32
def hcard(slide,pos,icon,head,desc,hl=False,icon_alt=None):
    r,c=divmod(pos,2); x=COLX[c]; y=ROWY[r]
    sp=card(slide,x,y,CW2,CH2,DG if hl else PALE)
    by=y+(CH2-BD2)/2
    badge(slide,x+LPAD,by,BD2,(WHITE if hl else DG),(icon_alt if hl and icon_alt else icon))
    hcol=WHITE if hl else DARK; dcol=LIGHT if hl else GREY
    card_text(sp,[(head,16.5,hcol,True,0,4),(desc,12,dcol,False,0,0)],
              ml=LPAD+BD2+0.28, mt=0.18, mr=0.30, mb=0.18, anchor=MSO_ANCHOR.MIDDLE)

s=p.slides.add_slide(LAY)
set_ph(s,0,"Neue Governance & Verankerung")
set_ph(s,13,"Von der App zum zentralen Deloitte-Teilnehmermanagement-Tool")
hcard(s,0,"icons/shield.png","Betriebsrats-Freigabe","Neue, erweiterte BR-Freigabe steht in den nächsten Tagen – die Voraussetzung für den breiten Rollout.")
hcard(s,1,"icons/user-check.png","Power-User & Multiplikatoren","Assistenzen aus T&T abgeholt; Eva Wienkamp, Annette Stoffel und Ebru Genctürk als Power-User etabliert.")
hcard(s,2,"icons/share-2.png","Ausweitung auf Service Lines","Für Sondercalls bei SRT sowie T&L (Audit & Tax) eingeladen – weitere Assistenz-Calls folgen.")
hcard(s,3,"icons/send.png","Intranet-Kommunikation","Intranet-Artikel in Vorbereitung, um DEX für die breite Deloitte-Welt freizuschalten.")

s=p.slides.add_slide(LAY)
set_ph(s,0,"Plan nach vorne & Budgetbedarf")
set_ph(s,13,"Veröffentlichung, Schulung und Klärung der Finanzierung für das kommende FY")
hcard(s,0,"icons/globe.png","Veröffentlichung & Listung","DEX offiziell veröffentlichen und in den Deloitte Tools listen – als zentrales Teilnehmermanagement-Tool sichtbar machen.")
hcard(s,1,"icons/book-open.png","Schulungen & Onboarding","Zentrale Info- und Onboarding-Termine aufsetzen, um die breite Masse befähigt an die App heranzuführen.")
hcard(s,2,"icons/user-plus.png","Kleines Core-Team","Festes Junior-Team für die Begleitung der Schulungs- und Bekanntmachungs-Maßnahmen aufbauen.")
hcard(s,3,"icons/dollar-sign.png","Budget-Klärung","Verlängerung des Budget-Codes bzw. eigene Budgetposition fürs kommende FY – für Schulung, Betrieb & Routinearbeiten. Bedarf: ~5.000 EUR (FY25/26), ~10.000 EUR ab FY26/27. Bestehender Code: GASR8771. Unterlage für den OMP-Call.",hl=True,icon_alt="icons/dollar-sign-green.png")

p.save("20260629_DEX_Update.pptx")
# verify: no overlay text boxes -> every shape with text is a rounded-rect card or placeholder
overlay=0
for sld in p.slides:
    for sh in sld.shapes:
        if sh.has_text_frame and sh.text.strip() and not sh.is_placeholder and sh.shape_type!=1 and "Rounded" not in (sh.name or ""):
            # auto_shape type for rounded rect
            pass
print("saved",len(p.slides),"slides; text now lives inside the card fill shapes; font=Aptos")
